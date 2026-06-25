from __future__ import annotations

import csv
from dataclasses import dataclass
from pathlib import Path

from .document_parsers import DocumentParsingError, ParserUnavailable, parse_with_mineru, should_use_mineru_for_pdf
from .document_parsers.media_parser import AUDIO_EXTENSIONS, VIDEO_EXTENSIONS, transcribe_audio, transcribe_video
from .document_parsers.office_converter import LEGACY_OFFICE_EXTENSIONS, convert_with_libreoffice
from .dependencies import detect_traditional_parser_dependencies


class UnsupportedDocumentType(ValueError):
    pass


@dataclass
class LoadedDocument:
    text: str
    metadata: dict[str, str]


TEXT_EXTENSIONS = {".txt", ".md", ".csv"}
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".tif", ".webp"}


def _read_text_file(path: Path) -> str:
    for encoding in ("utf-8-sig", "utf-8", "gb18030"):
        try:
            return path.read_text(encoding=encoding)
        except UnicodeDecodeError:
            continue
    return path.read_text(encoding="utf-8", errors="replace")


def _read_csv(path: Path) -> str:
    raw = _read_text_file(path)
    rows = []
    for row in csv.reader(raw.splitlines()):
        rows.append(",".join(str(cell).strip() for cell in row))
    return "\n".join(rows)


def _read_pdf(path: Path) -> tuple[str, int]:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    pages = []
    for index, page in enumerate(reader.pages):
        text = page.extract_text() or ""
        if text.strip():
            pages.append(f"[第 {index + 1} 页]\n{text.strip()}")
    return "\n\n".join(pages), len(reader.pages)


def _read_docx(path: Path) -> str:
    from docx import Document

    doc = Document(str(path))
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    table_rows = []
    for table in doc.tables:
        for row in table.rows:
            table_rows.append(" | ".join(cell.text.strip() for cell in row.cells))
    return "\n".join([*paragraphs, *table_rows])


def _read_xlsx(path: Path) -> str:
    from openpyxl import load_workbook

    workbook = load_workbook(str(path), read_only=True, data_only=True)
    lines = []
    for sheet in workbook.worksheets:
        lines.append(f"[Sheet: {sheet.title}]")
        for row in sheet.iter_rows(values_only=True):
            values = ["" if value is None else str(value).strip() for value in row]
            if any(values):
                lines.append(" | ".join(values))
    return "\n".join(lines)


def _read_pptx(path: Path) -> str:
    from pptx import Presentation

    presentation = Presentation(str(path))
    lines = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        lines.append(f"[Slide: {slide_index}]")
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text and text.strip():
                lines.append(text.strip())
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    values = [cell.text.strip() for cell in row.cells]
                    if any(values):
                        lines.append(" | ".join(values))
    return "\n".join(lines)


def _resolve_mineru_runtime_config() -> tuple[Path, str | list[str]]:
    import config

    output_root = Path(getattr(config, "RAGANYTHING_OUTPUT_ROOT", Path.cwd() / "output")) / "traditional_parser"
    command = getattr(config, "MINERU_COMMAND", None)
    if command:
        return output_root, list(command)
    mineru_path = str(getattr(config, "MINERU_CLI_PATH", "") or getattr(config, "MINERU_PATH", "") or "")
    return output_root, mineru_path


def _resolve_libreoffice_path() -> str:
    import config

    return str(getattr(config, "LIBREOFFICE_PATH", "") or "")


def _resolve_media_runtime_config() -> tuple[str, bool, Path]:
    import config

    deps = getattr(config, "TRADITIONAL_PARSER_DEPENDENCIES", {}) or {}
    current_deps = detect_traditional_parser_dependencies()
    ffmpeg_path = str(
        getattr(config, "FFMPEG_PATH", "")
        or deps.get("ffmpeg", {}).get("path", "")
        or current_deps.get("ffmpeg", {}).get("path", "")
        or ""
    )
    whisper_available = bool(
        getattr(config, "WHISPER_AVAILABLE", False)
        or deps.get("whisper", {}).get("available", False)
        or current_deps.get("whisper", {}).get("available", False)
    )
    output_root = Path(getattr(config, "RAGANYTHING_OUTPUT_ROOT", Path.cwd() / "output")) / "traditional_parser" / "media"
    return ffmpeg_path, whisper_available, output_root


def _load_with_mineru(file_path: Path) -> LoadedDocument:
    output_root, mineru_path = _resolve_mineru_runtime_config()
    parsed = parse_with_mineru(path=file_path, output_root=output_root, mineru_path=mineru_path)
    return LoadedDocument(
        text=parsed.text.strip(),
        metadata={
            "file_name": file_path.name,
            "extension": file_path.suffix.lower(),
            "path": str(file_path),
            **parsed.metadata,
        },
    )


def _is_mineru_unavailable(exc: Exception) -> bool:
    if isinstance(exc, ParserUnavailable):
        return True
    text = str(exc).lower()
    return "mineru cli not configured" in text or "mineru cli not found" in text


def _mineru_dependency_message(file_name: str) -> str:
    return (
        f"{file_name} 解析失败：该文件需要增强解析组件（MinerU）才能获得完整内容。"
        "请在桌面应用菜单中安装/修复增强解析组件后重试。"
    )


def _media_dependency_message(file_name: str, extension: str, exc: Exception) -> str:
    text = str(exc).lower()
    if extension in VIDEO_EXTENSIONS and "ffmpeg" in text:
        return (
            f"{file_name} 解析失败：视频解析需要 FFmpeg。"
            "请在桌面应用菜单中安装/修复增强解析组件后重试。"
        )
    if isinstance(exc, ParserUnavailable) and "whisper" in text:
        return (
            f"{file_name} 解析失败：音视频转写需要 Whisper 语音识别组件。"
            "当前基础版未安装该组件，请安装增强解析组件后重试。"
        )
    return f"{file_name} 解析失败：{exc}"


def load_document_text(path: str | Path) -> LoadedDocument:
    file_path = Path(path)
    extension = file_path.suffix.lower()
    if not file_path.exists() or not file_path.is_file():
        raise FileNotFoundError(str(file_path))

    if extension in {".txt", ".md"}:
        text = _read_text_file(file_path)
    elif extension == ".csv":
        text = _read_csv(file_path)
    elif extension == ".pdf":
        pdf_text, page_count = _read_pdf(file_path)
        if should_use_mineru_for_pdf(pdf_text, page_count):
            try:
                return _load_with_mineru(file_path)
            except (ParserUnavailable, DocumentParsingError) as exc:
                if _is_mineru_unavailable(exc) and str(pdf_text or "").strip():
                    text = pdf_text
                    return LoadedDocument(
                        text=text.strip(),
                        metadata={
                            "file_name": file_path.name,
                            "extension": extension,
                            "path": str(file_path),
                            "parser": "pypdf_fallback",
                            "warning": "增强解析组件（MinerU）不可用，已使用基础 PDF 文本解析结果。",
                        },
                    )
                if _is_mineru_unavailable(exc):
                    raise UnsupportedDocumentType(_mineru_dependency_message(file_path.name)) from exc
                raise UnsupportedDocumentType(f"{file_path.name} 解析失败：{exc}") from exc
        text = pdf_text
    elif extension == ".docx":
        text = _read_docx(file_path)
    elif extension == ".xlsx":
        text = _read_xlsx(file_path)
    elif extension == ".pptx":
        text = _read_pptx(file_path)
    elif extension in LEGACY_OFFICE_EXTENSIONS:
        try:
            converted_path = convert_with_libreoffice(
                path=file_path,
                output_dir=file_path.parent / "_converted_office",
                libreoffice_path=_resolve_libreoffice_path(),
            )
        except (ParserUnavailable, DocumentParsingError) as exc:
            raise UnsupportedDocumentType(f"{file_path.name} 解析失败：{exc}") from exc
        return load_document_text(converted_path)
    elif extension in IMAGE_EXTENSIONS:
        try:
            return _load_with_mineru(file_path)
        except (ParserUnavailable, DocumentParsingError) as exc:
            if _is_mineru_unavailable(exc):
                raise UnsupportedDocumentType(_mineru_dependency_message(file_path.name)) from exc
            raise UnsupportedDocumentType(f"{file_path.name} 解析失败：{exc}") from exc
    elif extension in AUDIO_EXTENSIONS:
        try:
            ffmpeg_path, whisper_available, _output_root = _resolve_media_runtime_config()
            text = transcribe_audio(file_path, whisper_available=whisper_available, ffmpeg_path=ffmpeg_path)
        except (ParserUnavailable, DocumentParsingError) as exc:
            raise UnsupportedDocumentType(_media_dependency_message(file_path.name, extension, exc)) from exc
    elif extension in VIDEO_EXTENSIONS:
        try:
            ffmpeg_path, whisper_available, output_root = _resolve_media_runtime_config()
            text = transcribe_video(
                file_path,
                output_dir=output_root / f"{file_path.stem}_audio",
                ffmpeg_path=ffmpeg_path,
                whisper_available=whisper_available,
            )
        except (ParserUnavailable, DocumentParsingError) as exc:
            raise UnsupportedDocumentType(_media_dependency_message(file_path.name, extension, exc)) from exc
    else:
        raise UnsupportedDocumentType(f"{file_path.name} 不支持传统 RAG 直接处理，请使用 RAG-Anything 高级解析。")

    return LoadedDocument(
        text=text.strip(),
        metadata={
            "file_name": file_path.name,
            "extension": extension,
            "path": str(file_path),
        },
    )
